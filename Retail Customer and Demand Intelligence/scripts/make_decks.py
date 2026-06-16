"""Build two landscape (16:9) decks from the project's real charts, aligned with README:
   Retailer-DS-Technical-Deck.pptx   (Senior Data Scientist voice)
   Retailer-CDO-Solution-Deck.pptx   (Chief Data Officer voice)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"

FOREST = RGBColor(0x21, 0x4E, 0x34)
MOSS   = RGBColor(0x6F, 0xA8, 0x7C)
CREAM  = RGBColor(0xF6, 0xF4, 0xEE)
CHAR   = RGBColor(0x2B, 0x2B, 0x2B)
GOLD   = RGBColor(0xC4, 0xA2, 0x65)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
RUST   = RGBColor(0xB8, 0x50, 0x42)
LIGHT  = RGBColor(0xD7, 0xE3, 0xDA)
EW, EH = Inches(13.333), Inches(7.5)

def deck():
    p = Presentation(); p.slide_width = EW; p.slide_height = EH; return p
def blank(p): return p.slides.add_slide(p.slide_layouts[6])
def rect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp
def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align; para.space_after = Pt(ln.get("sa", 4))
        r = para.add_run(); r.text = ln["t"]
        f = r.font; f.size = Pt(ln.get("sz", 16)); f.bold = ln.get("b", False)
        f.italic = ln.get("i", False); f.color.rgb = ln.get("c", CHAR); f.name = ln.get("fn", "Calibri")
    return tb
def img(s, path, l, t, maxw, maxh):
    w, h = Image.open(path).size; ar = w / h; bw, bh = maxw / Emu(1), maxh / Emu(1)
    if bw / bh > ar: dw = bh * ar; dh = bh
    else: dw = bw; dh = bw / ar
    pic = s.shapes.add_picture(str(path), 0, 0, width=Emu(int(dw)))
    pic.left = int(l + (maxw - pic.width) / 2); pic.top = int(t + (maxh - pic.height) / 2)
    return pic

def title_slide(p, kicker, title_lines, subtitle, footer):
    s = blank(p); rect(s, 0, 0, EW, EH, FOREST)
    rect(s, 0, Inches(6.7), EW, Inches(0.8), RGBColor(0x1A, 0x3D, 0x29))
    txt(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.5),
        [{"t": kicker, "sz": 18, "b": True, "c": GOLD, "fn": "Calibri"}])
    txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6),
        [{"t": t, "sz": 46, "b": True, "c": WHITE, "fn": "Georgia", "sa": 2} for t in title_lines])
    txt(s, Inches(0.9), Inches(4.7), Inches(11.0), Inches(1.2),
        [{"t": subtitle, "sz": 19, "c": LIGHT}])
    txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
        [{"t": footer, "sz": 12, "c": RGBColor(0xB9,0xCD,0xC0)}])
    return s

def header(s, kicker, title):
    rect(s, 0, 0, EW, EH, CREAM)
    rect(s, 0, 0, Inches(0.22), EH, FOREST)
    txt(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.4),
        [{"t": kicker, "sz": 13, "b": True, "c": MOSS}])
    txt(s, Inches(0.7), Inches(0.82), Inches(12), Inches(0.9),
        [{"t": title, "sz": 32, "b": True, "c": FOREST, "fn": "Georgia"}])

def stat_card(s, l, t, w, big, label, color=FOREST):
    rect(s, l, t, w, Inches(1.9), WHITE, line=RGBColor(0xE2,0xDE,0xD3))
    txt(s, l, t + Inches(0.22), w, Inches(1.0),
        [{"t": big, "sz": 34, "b": True, "c": color, "fn": "Georgia"}], align=PP_ALIGN.CENTER)
    txt(s, l + Inches(0.15), t + Inches(1.25), w - Inches(0.3), Inches(0.6),
        [{"t": label, "sz": 12.5, "c": GREY}], align=PP_ALIGN.CENTER)

def chart_slide(p, kicker, title, bullets, image, caption=""):
    s = blank(p); header(s, kicker, title)
    txt(s, Inches(0.7), Inches(1.95), Inches(4.5), Inches(4.8),
        [{"t": b, "sz": 15, "c": CHAR, "sa": 10} for b in bullets])
    img(s, image, Inches(5.4), Inches(1.9), Inches(7.4), Inches(4.7))
    if caption:
        txt(s, Inches(5.4), Inches(6.75), Inches(7.4), Inches(0.4),
            [{"t": caption, "sz": 10.5, "i": True, "c": GREY}])
    return s

def two_panel(p, kicker, title, left_head, left_lines, right_head, right_lines):
    s = blank(p); header(s, kicker, title)
    rect(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.6), RGBColor(0xF7,0xEC,0xEA), line=RGBColor(0xE6,0xD3,0xCF))
    txt(s, Inches(0.95), Inches(2.2), Inches(5.3), Inches(0.5), [{"t": left_head, "sz": 18, "b": True, "c": RUST}])
    txt(s, Inches(0.95), Inches(2.85), Inches(5.3), Inches(3.6), [{"t": x, "sz": 14, "c": CHAR, "sa": 8} for x in left_lines])
    rect(s, Inches(6.83), Inches(2.0), Inches(5.8), Inches(4.6), RGBColor(0xEC,0xF2,0xEE), line=RGBColor(0xCF,0xE0,0xD6))
    txt(s, Inches(7.08), Inches(2.2), Inches(5.3), Inches(0.5), [{"t": right_head, "sz": 18, "b": True, "c": FOREST}])
    txt(s, Inches(7.08), Inches(2.85), Inches(5.3), Inches(3.6), [{"t": x, "sz": 14, "c": CHAR, "sa": 8} for x in right_lines])
    return s

def bridge_slide(p, kicker, headline, body, flow):
    s = blank(p); rect(s, 0, 0, EW, EH, FOREST)
    txt(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5),
        [{"t": kicker, "sz": 16, "b": True, "c": GOLD}])
    txt(s, Inches(0.9), Inches(1.65), Inches(11.5), Inches(1.4),
        [{"t": headline, "sz": 38, "b": True, "c": WHITE, "fn": "Georgia"}])
    txt(s, Inches(0.9), Inches(3.45), Inches(11.4), Inches(2.4),
        [{"t": body, "sz": 18, "c": LIGHT}])
    rect(s, 0, Inches(6.25), EW, Inches(1.0), RGBColor(0x1A, 0x3D, 0x29))
    txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
        [{"t": flow, "sz": 15, "b": True, "c": GOLD}], anchor=MSO_ANCHOR.MIDDLE)
    return s

def closing_slide(p, title, lines, footer):
    s = blank(p); rect(s, 0, 0, EW, EH, FOREST)
    txt(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.0),
        [{"t": title, "sz": 40, "b": True, "c": WHITE, "fn": "Georgia"}])
    txt(s, Inches(0.9), Inches(3.0), Inches(11.2), Inches(3.0),
        [{"t": l, "sz": 18, "c": LIGHT, "sa": 10} for l in lines])
    txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
        [{"t": footer, "sz": 12, "c": RGBColor(0xB9,0xCD,0xC0)}])
    return s


def build_ds():
    p = deck()
    title_slide(p, "RETAIL CUSTOMER INTELLIGENCE  -  END-TO-END DATA SCIENCE",
                ["Customer & Demand", "Intelligence"],
                "Raw transactions -> cleaning -> star-schema warehouse -> 4 ML analyses -> profit-value model.  "
                "Business leads; technology serves.",
                "Senior Data Scientist portfolio - public UK Online Retail II as a proxy for a major retailer")

    two_panel(p, "PROBLEM & OBJECTIVE", "Why this project exists",
              "The problem",
              ["Every shopper treated the same",
               "Customers leave before anyone notices",
               "Promotions fired at everyone - discounting people who'd buy anyway",
               "Stock set by gut",
               "A chunk of revenue from shoppers we can't identify"],
              "The objective",
              ["Grow gross-margin profit by doing four things well:",
               "Know our customers; keep the valuable ones;",
               "grow the rest; stock the right products.",
               "Put a money value on each opportunity to rank investment.",
               "Prove every recommendation before spending at scale."])

    bridge_slide(p, "WHY THIS APPROACH",
                 "The answer is already in the data.",
                 "Every problem above is one gap - we don't truly know our customers. So the fix isn't more "
                 "discounts or more data; it's turning the transactions we already collect into four decisions, "
                 "each tied to profit. The method follows from that: business-led models, honest baselines, and a "
                 "controlled experiment behind every play - so we ship what's proven, not what's merely plausible.",
                 "Flying blind   ->   it's a customer-knowledge gap   ->   segment . predict churn . recommend . forecast")

    s = blank(p); header(s, "THE FOUNDATION", "From 1.07M messy rows to a trustworthy warehouse")
    stat_card(s, Inches(0.7), Inches(2.0), Inches(2.9), "94%", "of 1.07M rows kept (audited)")
    stat_card(s, Inches(3.8), Inches(2.0), Inches(2.9), "5 tables", "star schema (fact + 4 dims)")
    stat_card(s, Inches(6.9), Inches(2.0), Inches(2.9), "£19.64M", "revenue reconciled, 0 unmapped")
    stat_card(s, Inches(10.0), Inches(2.0), Inches(2.6), "12", "LLM-built product categories")
    txt(s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.4),
        [{"t": "Visible, reusable cleaning (duplicates, cancellations, non-positive, service codes); a modelling",
          "sz": 15, "sa": 4},
         {"t": "cleanse that keeps guests labelled Member/Non-Member; a feature foundation with the right", "sz": 15, "sa": 4},
         {"t": "source-of-truth split (customer models on Members; demand & basket guest-inclusive); and an", "sz": 15, "sa": 4},
         {"t": "LLM that created the product-category dimension the raw data lacked.", "sz": 15, "sa": 4}])

    chart_slide(p, "ANALYSIS 1 - SEGMENTATION", "Six behavioural segments (RFMTC + K-Means)",
                ["Features: Recency, Frequency, Monetary, Tenure + a churn-probability term (RFMTC).",
                 "k=6 chosen on elbow + silhouette + business capacity.",
                 "17% of members drive 70% of revenue.",
                 "Prioritised, one job each: Protect, Rescue, Grow."],
                FIG / "segment_matrix.png",
                "Six groups by value tier x lifecycle; priority segments colour-coded.")

    chart_slide(p, "ANALYSIS 1 - CLUSTERS", "Clusters in RFMTC space (PCA, 78% variance)",
                ["PCA projection of the standardised RFMTC features, coloured by segment.",
                 "Centroids marked; groups separate along a value-and-activity gradient.",
                 "Honest: silhouette ~0.36 - real groups, but a continuum, not islands."],
                FIG / "seg_cluster_scatter.png",
                "PCA = flatten the 5 RFMTC features onto a 2-D map so groups can be seen by eye; "
                "these two axes keep 78% of the real differences between customers.")

    chart_slide(p, "ANALYSIS 2 - CHURN", "Early-warning model - ranks who'll leave (AUC 0.79)",
                ["Strict time split: features before cutoff, label = no purchase in next 90 days.",
                 "AUC 0.79 = ranks a real leaver above a non-leaver ~4 times in 5 (0.5 = coin-flip).",
                 "Logistic ~ XGBoost - signal is mostly linear; reported honestly.",
                 "Top driver: recency vs the customer's normal rhythm. Feeds the rescue list."],
                FIG / "churn_roc.png", "ROC: both models ~0.79 vs the 0.5 chance line.")

    chart_slide(p, "ANALYSIS 2 - COHORTS", "Cohort retention - the first-repeat cliff",
                ["~20% return in month 1, settling to ~14% by month 6.",
                 "The leak is the SECOND purchase -> fix with onboarding, not loyalty tiers.",
                 "Cohort curves become the retention KPI to test every initiative.",
                 "Decay rate feeds lifetime value and the value model."],
                FIG / "cohort_retention_heatmap.png",
                "Curves are cohort averages; a few early cohorts tick UP near month 11-12 as seasonal "
                "(Christmas) buyers reactivate - retention isn't strictly monotonic in seasonal retail.")

    chart_slide(p, "ANALYSIS 3 - CROSS-SELL", "Recommender - right product ~1 in 4 (~100x chance)",
                ["Item-based collaborative filtering (cosine on the basket x item matrix).",
                 "Leave-one-out hit-rate@10 = 24.6% vs 0.24% random.",
                 "Affinity rules feed bundle ideas (with merchandiser curation).",
                 "Dead-stock paired with high-affinity anchors to clear inventory."],
                FIG / "cross_sell_top_bundles.png", "Top cross-sell bundles by lift.")

    chart_slide(p, "ANALYSIS 4 - FORECAST", "Demand forecast -> profit-aware stocking",
                ["Global XGBoost (log target) vs a 4-week moving-average baseline.",
                 "Honest finding: baseline wins aggregate; XGBoost wins 61% of SKUs -> HYBRID.",
                 "Real-unit error: MAE ~220 units; and it UNDER-forecasts ~126 units/wk -",
                 "left uncorrected that causes stock-outs, so the newsvendor rule adds the buffer."],
                FIG / "forecast_top_skus.png", "1-week-ahead forecast vs actual, top SKUs (12-week holdout).")

    chart_slide(p, "THE VALUE", "Levers priced in money - and a tangible campaign view",
                ["~£82k/yr incremental profit (+2.4% of profit); scales to ~£8M/yr at £1B revenue.",
                 "Retention is the biggest lever; ranges £41k-123k (uplift), £57k-106k (margin).",
                 "Rescue campaign: ~30 VIPs won back -> ~£43k revenue / ~£15k margin -> net ~£11k, ~3x return.",
                 "Every play ships with a holdout control to prove incremental lift."],
                FIG / "value_before_after.png", "Projected uplift; the same % scales with revenue.")

    s = blank(p); header(s, "ENGINEERING & HONESTY", "How it's built - and what it can't claim")
    txt(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(4.6),
        [{"t": "Stack", "sz": 18, "b": True, "c": FOREST, "sa": 6},
         {"t": "Python, pandas, scikit-learn, XGBoost, lifetimes (BG/NBD)", "sz": 14, "sa": 4},
         {"t": "scipy, matplotlib, parquet star schema, reproducible modules", "sz": 14, "sa": 4},
         {"t": "LLM-assisted category enrichment (GenAI); leakage-safe splits", "sz": 14, "sa": 12},
         {"t": "Evaluation", "sz": 18, "b": True, "c": FOREST, "sa": 6},
         {"t": "AUC for ranking; WAPE + MAE/bias (real units) for forecasts", "sz": 14, "sa": 4},
         {"t": "Always benchmarked vs a simple baseline; holdout controls", "sz": 14}])
    rect(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.6), WHITE, line=RGBColor(0xE2,0xDE,0xD3))
    txt(s, Inches(7.15), Inches(2.2), Inches(5.2), Inches(4.2),
        [{"t": "Honest caveats (stated, not hidden)", "sz": 16, "b": True, "c": CHAR, "sa": 8},
         {"t": "Public UK proxy - methods transfer, product mix does not.", "sz": 13.5, "sa": 6},
         {"t": "No COGS -> margin is an explicit, sensitivity-tested assumption.", "sz": 13.5, "sa": 6},
         {"t": "Lost sales unobservable -> availability lever is assumption-heavy.", "sz": 13.5, "sa": 6},
         {"t": "Small cohorts (rescue n=250) -> directional, expand to gain power.", "sz": 13.5, "sa": 6},
         {"t": "Simple beat ML on forecast aggregate -> hybrid, not a trophy model.", "sz": 13.5}])

    closing_slide(p, "From data to decisions",
                  ["Clean foundation -> warehouse -> segmentation, churn, cross-sell, forecasting -> profit value.",
                   "Business-first framing, leakage-safe modelling, honest metrics and baselines.",
                   "Code, methodology, and limitations: README.md and docs/."],
                  "Senior Data Scientist portfolio - retail-intelligence")
    out = ROOT / "Retailer-DS-Technical-Deck.pptx"; p.save(out); print("saved", out.name, "slides:", len(p.slides._sldIdLst))


def build_cdo():
    p = deck()
    title_slide(p, "OFFICE OF THE CHIEF DATA OFFICER",
                ["Grow margin by knowing,", "keeping & growing", "your best customers"],
                "A customer & demand intelligence capability that turns transaction data into profit - "
                "with every play proven by experiment.",
                "Retail customer-intelligence solution - outcome-focused")

    s = blank(p); header(s, "THE PROBLEM", "Before: flying without instruments")
    rows = [("No customer view", "every shopper treated alike"),
            ("Churn found too late", "noticed only after they've gone"),
            ("Marketing by gut", "blanket discounts erode margin"),
            ("Generic cross-sell", "no 'what to offer next'"),
            ("Reactive inventory", "stockouts lose sale + basket"),
            ("13% revenue invisible", "guest sales untracked")]
    for i, (a, b) in enumerate(rows):
        c, r = i % 2, i // 2
        l = Inches(0.7 + c * 6.2); t = Inches(2.1 + r * 1.5)
        rect(s, l, t, Inches(5.9), Inches(1.25), WHITE, line=RGBColor(0xE2,0xDE,0xD3))
        txt(s, l + Inches(0.25), t + Inches(0.18), Inches(5.4), Inches(0.9),
            [{"t": a, "sz": 16, "b": True, "c": RUST, "sa": 2},
             {"t": b, "sz": 13, "c": GREY}])

    s = blank(p); header(s, "THE OBJECTIVE", "Grow gross-margin profit - four jobs")
    obj = [("KNOW", "Who our customers are", "segmentation"),
           ("KEEP", "Retain the valuable ones", "churn + win-back"),
           ("GROW", "Lift everyone else", "cross-sell"),
           ("STOCK", "The right products", "demand forecast")]
    for i, (k, line, meth) in enumerate(obj):
        l = Inches(0.7 + i * 3.07)
        rect(s, l, Inches(2.2), Inches(2.85), Inches(3.1), WHITE, line=RGBColor(0xE2,0xDE,0xD3))
        rect(s, l, Inches(2.2), Inches(2.85), Inches(0.65), FOREST)
        txt(s, l, Inches(2.31), Inches(2.85), Inches(0.45), [{"t": k, "sz": 16, "b": True, "c": WHITE}], align=PP_ALIGN.CENTER)
        txt(s, l + Inches(0.2), Inches(3.15), Inches(2.45), Inches(1.4),
            [{"t": line, "sz": 16, "b": True, "c": FOREST, "fn": "Georgia"}], align=PP_ALIGN.CENTER)
        txt(s, l + Inches(0.2), Inches(4.55), Inches(2.45), Inches(0.5),
            [{"t": "(" + meth + ")", "sz": 12, "i": True, "c": GREY}], align=PP_ALIGN.CENTER)
    rect(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.8), RGBColor(0xEC,0xF2,0xEE), line=RGBColor(0xCF,0xE0,0xD6))
    txt(s, Inches(0.95), Inches(5.92), Inches(11.4), Inches(0.5),
        [{"t": "Each opportunity priced in money to rank investment; each play proven by a controlled "
               "experiment before we scale.", "sz": 14, "b": True, "c": FOREST}])

    bridge_slide(p, "WHY THIS WORKS",
                 "The answer is already in your data.",
                 "Every problem is the same gap - you don't yet know your customers. The fix isn't more discounts "
                 "or new systems; it's turning the transactions you already collect into four profit decisions, and "
                 "proving each one with a controlled test before you spend at scale. Low risk, measured return.",
                 "Flying blind   ->   a customer-knowledge gap   ->   know . keep . grow . stock")

    s = blank(p); header(s, "THE SOLUTION", "One capability, three questions answered")
    steps = [("WHO", "Segment customers", "6 behavioural groups; 17% drive 70% of revenue"),
             ("WHO'S LEAVING", "Spot leavers early", "flagged before they go, ranked accurately"),
             ("WHAT TO OFFER", "Recommend next product", "personalised cross-sell, ~100x better than random")]
    for i, (k, h, d) in enumerate(steps):
        l = Inches(0.7 + i * 4.15)
        rect(s, l, Inches(2.2), Inches(3.85), Inches(3.4), WHITE, line=RGBColor(0xE2,0xDE,0xD3))
        rect(s, l, Inches(2.2), Inches(3.85), Inches(0.7), FOREST)
        txt(s, l, Inches(2.32), Inches(3.85), Inches(0.5), [{"t": k, "sz": 15, "b": True, "c": WHITE}], align=PP_ALIGN.CENTER)
        txt(s, l + Inches(0.3), Inches(3.15), Inches(3.25), Inches(0.7), [{"t": h, "sz": 19, "b": True, "c": FOREST, "fn": "Georgia"}])
        txt(s, l + Inches(0.3), Inches(4.0), Inches(3.25), Inches(1.4), [{"t": d, "sz": 14, "c": CHAR}])
    txt(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.6),
        [{"t": "+ a demand forecast that sets profit-optimal stock, and a value model that ranks every lever.",
          "sz": 14, "i": True, "c": GREY}])

    s = blank(p); header(s, "BEFORE  ->  AFTER", "What changes when data leads")
    pairs = [("Customer understanding", "One-size-fits-all", "6 segments, prioritised by profit"),
             ("Churn", "Found after they leave", "Scored early - 250-name, £0.7M rescue list"),
             ("Marketing spend", "Blanket, gut-led", "Profit-thresholded + proven by control group"),
             ("Cross-sell", "Generic", "Personalised next-best product"),
             ("Inventory", "Reactive", "Forecast + margin-aware stocking"),
             ("Guest revenue", "13% invisible (£2.6M)", "Sized & converted via loyalty")]
    txt(s, Inches(0.7), Inches(1.95), Inches(4.0), Inches(0.4), [{"t": "CAPABILITY", "sz": 12, "b": True, "c": GREY}])
    txt(s, Inches(4.8), Inches(1.95), Inches(3.7), Inches(0.4), [{"t": "BEFORE", "sz": 12, "b": True, "c": RUST}])
    txt(s, Inches(8.7), Inches(1.95), Inches(4.2), Inches(0.4), [{"t": "AFTER", "sz": 12, "b": True, "c": FOREST}])
    for i, (cap, bef, aft) in enumerate(pairs):
        t = Inches(2.4 + i * 0.78)
        if i % 2 == 0: rect(s, Inches(0.6), t, Inches(12.2), Inches(0.72), RGBColor(0xEE,0xEB,0xE2))
        txt(s, Inches(0.7), t + Inches(0.1), Inches(4.0), Inches(0.6), [{"t": cap, "sz": 13.5, "b": True, "c": CHAR}])
        txt(s, Inches(4.8), t + Inches(0.1), Inches(3.7), Inches(0.6), [{"t": bef, "sz": 13, "c": GREY}])
        txt(s, Inches(8.7), t + Inches(0.1), Inches(4.2), Inches(0.6), [{"t": aft, "sz": 13, "b": True, "c": FOREST}])

    chart_slide(p, "THE PLAYS", "Where the money is: Protect, Rescue, Grow",
                ["PROTECT - 70% of revenue, ~2% churn. Defend with loyalty & service, not discounts.",
                 "RESCUE - 250 high-value lapsing, ~£0.7M at risk. One-time win-back, proven by control.",
                 "GROW - 1,545 mid-tier with headroom. Cross-sell to lift spend, only where it's profitable.",
                 "Loyalty also converts the £2.6M anonymous guest base into tracked members."],
                FIG / "segment_matrix.png", "Six groups by value tier x lifecycle; priorities colour-coded.")

    chart_slide(p, "THE OUTCOME", "Proven profit - and it scales",
                ["~£82k/yr extra profit here (+2.4% of profit); scales to ~£8M/yr at £1B revenue.",
                 "Rescue campaign: ~30 customers won back -> net ~£11k, about a 3x return.",
                 "Bigger prize still: ~£0.7M at-risk + £2.6M untracked guest revenue to capture.",
                 "Retention is the biggest lever; each play proven by a holdout control before scaling."],
                FIG / "value_before_after.png", "Projected uplift; conservative, sensitivity-tested.")

    s = blank(p); header(s, "HOW I WORK", "Honest, business-first, evidence-led")
    txt(s, Inches(0.7), Inches(2.1), Inches(12), Inches(4.4),
        [{"t": "Business leads, technology serves - every model maps to a profit decision.", "sz": 17, "sa": 12},
         {"t": "Money on the metrics - results in money, not just model scores.", "sz": 17, "sa": 12},
         {"t": "Proven, not promised - controlled experiments before spending at scale.", "sz": 17, "sa": 12},
         {"t": "Candid about limits - assumptions stated, simple baselines respected.", "sz": 17, "sa": 12},
         {"t": "Reproducible - clean code, a warehouse, and documentation that doesn't drift.", "sz": 17}])

    closing_slide(p, "What I can do for you",
                  ["Stand up customer segmentation, churn early-warning, and cross-sell engines.",
                   "Turn your transaction data into ranked, costed, experiment-backed growth plays.",
                   "Deliver it reproducibly - warehouse, models, and an executive narrative leaders trust."],
                  "Let's talk about your data -> profit roadmap.")
    out = ROOT / "Retailer-CDO-Solution-Deck.pptx"; p.save(out); print("saved", out.name, "slides:", len(p.slides._sldIdLst))


if __name__ == "__main__":
    build_ds(); build_cdo()
